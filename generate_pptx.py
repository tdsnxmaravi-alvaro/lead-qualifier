from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x4E, 0xC9, 0xB0)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)


def add_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_bullet_frame(slide, left, top, width, height, items, font_size=13, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (label, text) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        if label:
            run = p.add_run()
            run.text = label
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = GREEN
            run = p.add_run()
            run.text = f"  {text}"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return tf


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_bg(slide1)

add_text_box(slide1, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
             "AI-Powered Lead Qualification Engine", font_size=36, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide1, Inches(1), Inches(3.4), Inches(11), Inches(0.8),
             "Copilot Studio Multi-Agent System for TD SYNNEX Datech", font_size=20,
             color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide1, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "Pilot Results & Scale Roadmap  |  May 2026", font_size=16,
             color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Architecture line
add_text_box(slide1, Inches(2.5), Inches(5.5), Inches(8), Inches(1.2),
             "Orchestrator  \u2192  Newforma Agent  |  Bricsys Agent  |  + extensible to Unity, DraftSight, Novade",
             font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: EXECUTIVE SUMMARY
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide2)

add_text_box(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "Executive Summary", font_size=26, bold=True, color=WHITE)

# Left column
left_items = [
    ("CHALLENGE:", "No systematic way to qualify 2,000+ resellers for vendor-specific fit. Manual research = weeks per batch, anecdotal, unscalable."),
    ("GOAL:", "Evidence-backed, AI-driven reseller scoring \u2014 BDM-ready intelligence per account with product recommendations and talking points."),
    ("WHAT WE DID:", "Built a multi-agent system in Copilot Studio: one orchestrator routes to specialized agents (Newforma, Bricsys). Each agent researches, scores, and ranks resellers using vendor-specific criteria via web search + structured prompts."),
]

add_bullet_frame(slide2, Inches(0.5), Inches(1.0), Inches(6.0), Inches(5.5),
                 left_items, font_size=13)

# Right column
right_items = [
    ("OUTCOME:", "Pilot validated \u2014 BDM confirmed AI-generated intel matched real customer data. Quote: \"Had feedback from lead tool, one was for a reseller\u2019s current customer. He confirmed all info was correct.\""),
    ("KEY TAKEAWAY:", "The system works. Accurate, scalable, and ready to extend to all 5 Datech vendors (Newforma, Bricsys, Unity, DraftSight, Novade)."),
]

add_bullet_frame(slide2, Inches(6.8), Inches(1.0), Inches(6.0), Inches(5.5),
                 right_items, font_size=13)


# ============================================================
# SLIDE 3: ROI & SCALE
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide3)

add_text_box(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             "ROI & Measurement", font_size=26, bold=True, color=WHITE)

# ROI table as text block
roi_items = [
    ("AI Pipeline:", "50 resellers qualified in ~2.5 hrs / ~$50"),
    ("Manual Equivalent:", "25\u201350 hrs / $1,875\u2013$3,750 per batch"),
    ("Scale:", "Same marginal cost for 2,000+ resellers vs. 6+ months BDM calendar manually"),
    ("Output:", "Scored profiles + product recs + talking points + cited sources (vs. unstructured CRM notes)"),
]

add_bullet_frame(slide3, Inches(0.5), Inches(1.0), Inches(5.8), Inches(3.5),
                 roi_items, font_size=13)

# Measurement (To-Be)
add_text_box(slide3, Inches(7.0), Inches(0.9), Inches(5.5), Inches(0.5),
             "Measurement (To-Be)", font_size=16, bold=True, color=GREEN)

measure_items = [
    ("", "\u2022 # resellers recruited with active pipeline within 90 days"),
    ("", "\u2022 Time from qualification \u2192 first BDM engagement"),
    ("", "\u2022 Revenue attributed to AI-identified accounts"),
    ("", "\u2022 Vendor MDF secured for targeted campaigns"),
]

add_bullet_frame(slide3, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.5),
                 measure_items, font_size=12, color=LIGHT_GRAY)

# Next Steps
add_text_box(slide3, Inches(0.5), Inches(5.0), Inches(12), Inches(0.5),
             "Next Steps", font_size=16, bold=True, color=ORANGE)

next_items = [
    ("", "1. Scale Newforma & Bricsys agents to full 2,000+ reseller base"),
    ("", "2. Deploy new agents for Unity, DraftSight, Novade (same architecture, new prompts)"),
    ("", "3. Integrate output into CRM workflow for automated BDM assignment"),
]

add_bullet_frame(slide3, Inches(0.5), Inches(5.5), Inches(12), Inches(1.8),
                 next_items, font_size=13)

# Save
prs.save(r"c:\dev\lead-qualifier\Lead_Qualifier_Exec_Presentation.pptx")
print("Presentation saved: Lead_Qualifier_Exec_Presentation.pptx")
